from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "Selenium_Automation_Framework_AI_Healing.pptx"


# --- Theme colors (bright but professional) ---
INK = RGBColor(234, 242, 255)
MUTED = RGBColor(190, 205, 230)
BG = RGBColor(11, 16, 32)

C1 = RGBColor(110, 231, 255)  # cyan
C2 = RGBColor(167, 139, 250)  # purple
C3 = RGBColor(52, 211, 153)   # green
C4 = RGBColor(245, 158, 11)   # amber


def _set_text(run, text: str, size: int, bold: bool = False, color=INK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    # Title block
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.6), Inches(12.0), Inches(1.4)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(2, 6, 23)  # deep navy
    box.fill.transparency = 0.25
    box.line.color.rgb = RGBColor(255, 255, 255)
    box.line.transparency = 0.85
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    _set_text(r, title, 34, True, INK)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        _set_text(r2, subtitle, 16, False, MUTED)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    _set_text(r, text, 10, False, MUTED)


def add_bullets(slide, x, y, w, h, items: list[str], font_size=20):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK
    return box


def kpi_card(slide, x, y, w, h, heading: str, value: str, accent: RGBColor):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(2, 6, 23)
    card.fill.transparency = 0.18
    card.line.color.rgb = RGBColor(255, 255, 255)
    card.line.transparency = 0.86

    # Accent bar
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.10), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    _set_text(r1, heading, 12, False, MUTED)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    _set_text(r2, value, 18, True, INK)


def arrow(slide, x1, y1, x2, y2, color=MUTED):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    return conn


def flow_box(slide, x, y, w, h, text: str, accent: RGBColor):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(2, 6, 23)
    box.fill.transparency = 0.18
    box.line.color.rgb = accent
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _set_text(r, text, 16, True, INK)
    return box


def make():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    cover = ASSETS / "cover.png"
    arch = ASSETS / "architecture.png"
    heal = ASSETS / "ai-healing.png"

    # --- Slide 1: Cover ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    s.shapes.add_picture(str(cover), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    add_title(s, "Selenium Automation Framework", "Java 17 • TestNG • Smart Actions • AI Auto‑Healing (Optional)")
    add_footer(s, "Automation Assignment • Framework Demo")

    # --- Slide 2: What it delivers ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "What this framework delivers", "Maintainability • Stability • Resilience")
    kpi_card(s, Inches(0.9), Inches(2.2), Inches(3.9), Inches(1.2), "Stability", "Explicit waits", C1)
    kpi_card(s, Inches(0.9), Inches(3.55), Inches(3.9), Inches(1.2), "Maintainability", "POM + annotations", C2)
    kpi_card(s, Inches(0.9), Inches(4.9), Inches(3.9), Inches(1.2), "Resilience", "AI healing + audit", C3)
    add_bullets(
        s,
        Inches(0.9),
        Inches(6.2),
        Inches(6.2),
        Inches(0.7),
        [
            "Low setup with Selenium Manager",
            "Centralized waits + smart actions",
            "Safety: redaction, cache, validation",
        ],
        font_size=16,
    )
    s.shapes.add_picture(str(arch), Inches(7.1), Inches(2.2), width=Inches(5.5))
    add_footer(s, "Slide 2/10")

    # --- Slide 3: Problem statement (diagram) ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "The problem we solve", "UI changes cause fragile selectors and maintenance cost")
    b1 = flow_box(s, Inches(1.0), Inches(2.7), Inches(2.6), Inches(0.9), "UI Changes", C4)
    b2 = flow_box(s, Inches(4.0), Inches(2.7), Inches(2.6), Inches(0.9), "Broken Selectors", C2)
    b3 = flow_box(s, Inches(7.0), Inches(2.7), Inches(2.6), Inches(0.9), "Test Failures", C1)
    b4 = flow_box(s, Inches(10.0), Inches(2.7), Inches(2.6), Inches(0.9), "Maintenance Cost", C3)
    arrow(s, b1.left + b1.width, b1.top + b1.height / 2, b2.left, b2.top + b2.height / 2)
    arrow(s, b2.left + b2.width, b2.top + b2.height / 2, b3.left, b3.top + b3.height / 2)
    arrow(s, b3.left + b3.width, b3.top + b3.height / 2, b4.left, b4.top + b4.height / 2)
    add_bullets(
        s,
        Inches(1.0),
        Inches(4.3),
        Inches(11.6),
        Inches(2.6),
        [
            "Flaky runs due to timing and dynamic DOM",
            "Selector churn during UI releases",
            "Boilerplate waits and locators repeated across tests",
            "Limited traceability on what changed and why",
        ],
        font_size=18,
    )
    add_footer(s, "Slide 3/10")

    # --- Slide 4: Architecture ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "Architecture overview", "Layered design with clear responsibilities")
    s.shapes.add_picture(str(arch), Inches(0.9), Inches(2.2), width=Inches(6.1))
    add_bullets(
        s,
        Inches(7.3),
        Inches(2.3),
        Inches(5.1),
        Inches(4.8),
        [
            "Tests (TestNG): call page actions",
            "Pages: business actions + @SeleniumSelector",
            "Core: waits + smart actions + element wiring",
            "AI (optional): healing with audit log",
            "Utilities: JSON + HTML cleaning",
        ],
        font_size=18,
    )
    add_footer(s, "Slide 4/10")

    # --- Slide 5: Core building blocks ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "Core implementation", "How the framework works under the hood")
    cA = flow_box(s, Inches(1.0), Inches(2.4), Inches(3.7), Inches(1.0), "SeleniumFactory\n(WebDriver lifecycle)", C1)
    cB = flow_box(s, Inches(1.0), Inches(3.8), Inches(3.7), Inches(1.0), "BasePage\n(waits + smart actions)", C2)
    cC = flow_box(s, Inches(5.2), Inches(2.4), Inches(3.7), Inches(1.0), "ElementInitializer\n(reflection wiring)", C3)
    cD = flow_box(s, Inches(5.2), Inches(3.8), Inches(3.7), Inches(1.0), "SmartElement\n(By + driver)", C4)
    arrow(s, cA.left + cA.width, cA.top + cA.height / 2, cC.left, cC.top + cC.height / 2, color=MUTED)
    arrow(s, cB.left + cB.width, cB.top + cB.height / 2, cD.left, cD.top + cD.height / 2, color=MUTED)
    add_bullets(
        s,
        Inches(9.3),
        Inches(2.4),
        Inches(3.9),
        Inches(4.2),
        [
            "Explicit waits (10s default)",
            "CSS selectors by default",
            "XPath via prefix: xpath=...",
            "Minimal boilerplate in tests",
        ],
        font_size=18,
    )
    add_footer(s, "Slide 5/10")

    # --- Slide 6: Smart actions ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "Smart Actions", "Stable by default; healing only on failure")
    add_bullets(
        s,
        Inches(0.9),
        Inches(2.2),
        Inches(6.3),
        Inches(4.8),
        [
            "smartClick(selector): wait → click",
            "smartFill(selector, value): wait → clear → type",
            "On TimeoutException → healing attempt (1 retry)",
            "StaleElementReferenceException → retry once (no AI)",
        ],
        font_size=20,
    )
    # Visual ladder
    step1 = flow_box(s, Inches(7.4), Inches(2.3), Inches(5.2), Inches(0.75), "Parse selector → By", C1)
    step2 = flow_box(s, Inches(7.4), Inches(3.2), Inches(5.2), Inches(0.75), "Explicit wait (visibility/clickable)", C2)
    step3 = flow_box(s, Inches(7.4), Inches(4.1), Inches(5.2), Inches(0.75), "Perform action (click/type)", C3)
    step4 = flow_box(s, Inches(7.4), Inches(5.0), Inches(5.2), Inches(0.75), "If timeout → AI heal → validate → retry once", C4)
    arrow(s, step1.left + step1.width / 2, step1.top + step1.height, step2.left + step2.width / 2, step2.top)
    arrow(s, step2.left + step2.width / 2, step2.top + step2.height, step3.left + step3.width / 2, step3.top)
    arrow(s, step3.left + step3.width / 2, step3.top + step3.height, step4.left + step4.width / 2, step4.top)
    add_footer(s, "Slide 6/10")

    # --- Slide 7: AI healing pipeline ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "AI Healing pipeline", "Optional • audited • designed with safety guardrails")
    s.shapes.add_picture(str(heal), Inches(0.9), Inches(2.2), width=Inches(6.2))
    add_bullets(
        s,
        Inches(7.4),
        Inches(2.3),
        Inches(5.2),
        Inches(4.9),
        [
            "Trigger: action times out (TimeoutException)",
            "Capture HTML via getPageSource()",
            "Clean + redact sensitive values (best-effort)",
            "Ask LLM for replacement selector (if configured)",
            "Validate healed selector → retry once",
            "Audit log: logs/ai-healing-audit.log",
        ],
        font_size=18,
    )
    add_footer(s, "Slide 7/10")

    # --- Slide 8: Hardening: cache + validation ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "AI hardening (real projects)", "Redaction • Healing cache • Selector validation")
    hb1 = flow_box(s, Inches(1.0), Inches(2.5), Inches(3.6), Inches(0.9), "Cache lookup\nhealing-cache.json", C3)
    hb2 = flow_box(s, Inches(5.0), Inches(2.5), Inches(3.6), Inches(0.9), "Validate selector\n(unique + actionable)", C1)
    hb3 = flow_box(s, Inches(9.0), Inches(2.5), Inches(3.6), Inches(0.9), "Retry once\n(or bypass cache)", C4)
    arrow(s, hb1.left + hb1.width, hb1.top + hb1.height / 2, hb2.left, hb2.top + hb2.height / 2)
    arrow(s, hb2.left + hb2.width, hb2.top + hb2.height / 2, hb3.left, hb3.top + hb3.height / 2)
    add_bullets(
        s,
        Inches(1.0),
        Inches(3.8),
        Inches(11.6),
        Inches(3.0),
        [
            "Cache reduces repeated LLM calls and speeds up healing",
            "Validation prevents wrong/stale selectors from hiding real issues",
            "If cached selector is invalid: bypass cache once and request fresh heal",
            "Override cache path: mvn test -Dhealing.cache.path=...json",
        ],
        font_size=18,
    )
    add_footer(s, "Slide 8/10")

    # --- Slide 9: Demo story ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "Demo story", "Login test flow + evidence artifacts")
    add_bullets(
        s,
        Inches(0.9),
        Inches(2.2),
        Inches(7.0),
        Inches(4.8),
        [
            "Test creates driver via SeleniumFactory",
            "LoginPage exposes clear actions (enterUsername, enterPassword, clickLogin)",
            "Smart actions add stability and optional healing on failures",
            "Evidence: ai-healing-audit.log + healing-cache.json",
        ],
        font_size=20,
    )
    cmd = flow_box(s, Inches(8.2), Inches(2.6), Inches(4.4), Inches(0.85), "Run: mvn test", C2)
    evid1 = flow_box(s, Inches(8.2), Inches(3.7), Inches(4.4), Inches(0.85), "Audit: logs/ai-healing-audit.log", C4)
    evid2 = flow_box(s, Inches(8.2), Inches(4.8), Inches(4.4), Inches(0.85), "Cache: healing-cache.json", C3)
    add_footer(s, "Slide 9/10")

    # --- Slide 10: Summary / Why ---
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_title(s, "Why this framework", "Clean • Reliable • Explainable • Resilient")
    kpi_card(s, Inches(0.9), Inches(2.3), Inches(4.0), Inches(1.3), "Clarity", "Page Object Model", C2)
    kpi_card(s, Inches(0.9), Inches(3.75), Inches(4.0), Inches(1.3), "Reliability", "Explicit waits", C1)
    kpi_card(s, Inches(0.9), Inches(5.2), Inches(4.0), Inches(1.3), "Resilience", "AI healing (optional)", C3)
    add_bullets(
        s,
        Inches(5.3),
        Inches(2.3),
        Inches(7.0),
        Inches(4.6),
        [
            "Designed to reduce flakiness and locator maintenance",
            "AI hardening: redaction + cache + validation guardrails",
            "Audit trail provides traceability for assignments and real teams",
            "Future scope: ThreadLocal drivers, reporting, screenshots",
        ],
        font_size=20,
    )
    end = flow_box(s, Inches(5.3), Inches(6.3), Inches(7.0), Inches(0.75), "Thank you — Questions?", C4)
    add_footer(s, "Slide 10/10")

    prs.save(str(OUT))
    print(f"Created: {OUT}")


if __name__ == "__main__":
    make()

